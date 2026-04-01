"""
crop_element_thumbnails.py
从 LibreOffice 渲染的页面图片中，根据元素坐标裁剪出每个素材的真实缩略图。

流程：
1. 读取 assetIndex.json 获取所有素材及其 slideNumber + elementIndex
2. 用 python-pptx 读取源 PPTX，获取每个元素的精确坐标 (EMU)
3. 从渲染好的页面 PNG 中裁剪对应区域
4. 保存为缩略图

注意：PPTX 的 slide 编号和 PDF 的页码之间可能不对应（PPTX slide 可能不是连续编号的），
所以需要用 python-pptx 按顺序读取幻灯片，建立 slide 到页码的映射。
"""

import json
import os
import sys
import xml.etree.ElementTree as ET
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

# ---- Config ----
PROJECT_DIR = Path(__file__).parent.parent
PPTX_FILE = None  # Will be found via glob
PAGES_DIR = Path("/tmp/pptx_render/pages")
THUMB_DIR = PROJECT_DIR / "assets" / "thumbnails"
INDEX_FILE = PROJECT_DIR / "src" / "data" / "assetIndex.json"
THUMB_SIZE = (300, 225)  # Target thumbnail size

# 渲染参数（必须和渲染时一致）
ZOOM = 2.0  # PDF 渲染缩放倍数

# PPTX 标准尺寸 (EMU)
# python-pptx 会给出精确的 slide 尺寸

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def find_pptx():
    """找到源 PPTX 文件"""
    import glob
    files = glob.glob(str(PROJECT_DIR / "*.pptx"))
    if not files:
        print("❌ No .pptx file found in project root!")
        sys.exit(1)
    return files[0]


def get_slide_number_mapping(unpacked_dir):
    """
    建立 slide 序号到物理文件编号的映射。
    PPTX 的 presentation.xml 中按顺序列出了所有 slide 的关系 ID，
    但物理文件名可能不连续（如 slide1.xml, slide5.xml, slide3.xml...）
    
    返回: {0-based index: physical slide number}
    """
    pres_xml = unpacked_dir / "ppt" / "presentation.xml"
    pres_rels = unpacked_dir / "ppt" / "_rels" / "presentation.xml.rels"
    
    if not pres_xml.exists() or not pres_rels.exists():
        return None
    
    # 读取关系文件，建立 rId -> target 的映射
    rels_tree = ET.parse(pres_rels)
    rels_root = rels_tree.getroot()
    rid_to_target = {}
    for rel in rels_root.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
        rid = rel.get('Id', '')
        target = rel.get('Target', '')
        if 'slide' in target.lower() and 'slideLayout' not in target and 'slideMaster' not in target:
            rid_to_target[rid] = target
    
    # 读取 presentation.xml 中的 slide 引用顺序
    pres_tree = ET.parse(pres_xml)
    pres_root = pres_tree.getroot()
    
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    
    slide_list = pres_root.find(f'{{{ns_p}}}sldIdLst')
    if slide_list is None:
        return None
    
    ordered_slides = []
    for sld_id in slide_list:
        rid = sld_id.get(f'{{{ns_r}}}id', '')
        target = rid_to_target.get(rid, '')
        # 提取 slide 编号
        import re
        m = re.search(r'slide(\d+)\.xml', target)
        if m:
            ordered_slides.append(int(m.group(1)))
    
    # mapping: page_index (0-based) -> slide physical number
    mapping = {}
    for idx, snum in enumerate(ordered_slides):
        mapping[snum] = idx  # slide number -> page index (0-based)
    
    return mapping


def get_element_bounds_from_pptx(pptx_path):
    """
    用 python-pptx 读取每张幻灯片上每个元素的位置和尺寸。
    返回: {slide_index(0-based): [(element_name, left, top, width, height), ...]}
    
    坐标单位为 EMU (English Metric Units)
    """
    prs = Presentation(pptx_path)
    slide_width = prs.slide_width  # EMU
    slide_height = prs.slide_height  # EMU
    
    print(f"Slide dimensions: {slide_width}x{slide_height} EMU "
          f"({Emu(slide_width).inches:.2f}x{Emu(slide_height).inches:.2f} inches)")
    
    all_bounds = {}
    
    for slide_idx, slide in enumerate(prs.slides):
        elements = []
        for shape_idx, shape in enumerate(slide.shapes):
            # 跳过 placeholder（标题等固定元素）
            if shape.is_placeholder:
                continue
            
            elements.append({
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'shape_type': str(shape.shape_type) if hasattr(shape, 'shape_type') else 'unknown',
            })
        
        all_bounds[slide_idx] = elements
    
    return all_bounds, slide_width, slide_height, len(prs.slides)


def crop_element_thumbnail(page_img, element_bounds, slide_width_emu, slide_height_emu, padding_pct=0.05):
    """
    从页面图片中裁剪出元素区域。
    
    Args:
        page_img: PIL Image of the full page
        element_bounds: dict with left, top, width, height in EMU
        slide_width_emu, slide_height_emu: slide dimensions in EMU
        padding_pct: padding as percentage of element size
    
    Returns: PIL Image of the cropped element
    """
    img_w, img_h = page_img.size
    
    # EMU 到图片像素的转换比例
    scale_x = img_w / slide_width_emu
    scale_y = img_h / slide_height_emu
    
    # 转换坐标
    left = element_bounds['left'] * scale_x
    top = element_bounds['top'] * scale_y
    width = element_bounds['width'] * scale_x
    height = element_bounds['height'] * scale_y
    
    # 添加 padding
    pad_x = width * padding_pct
    pad_y = height * padding_pct
    
    x1 = max(0, left - pad_x)
    y1 = max(0, top - pad_y)
    x2 = min(img_w, left + width + pad_x)
    y2 = min(img_h, top + height + pad_y)
    
    # 裁剪
    cropped = page_img.crop((int(x1), int(y1), int(x2), int(y2)))
    
    # 缩放到目标尺寸（保持比例）
    thumb = Image.new('RGB', THUMB_SIZE, (255, 255, 255))
    cropped.thumbnail((THUMB_SIZE[0] - 8, THUMB_SIZE[1] - 8), Image.LANCZOS)
    # 居中放置
    paste_x = (THUMB_SIZE[0] - cropped.width) // 2
    paste_y = (THUMB_SIZE[1] - cropped.height) // 2
    
    if cropped.mode == 'RGBA':
        thumb.paste(cropped, (paste_x, paste_y), cropped)
    else:
        thumb.paste(cropped, (paste_x, paste_y))
    
    return thumb


def render_full_slide_thumbnail(page_img):
    """将整页渲染图缩放为缩略图"""
    thumb = Image.new('RGB', THUMB_SIZE, (255, 255, 255))
    page_copy = page_img.copy()
    page_copy.thumbnail((THUMB_SIZE[0] - 4, THUMB_SIZE[1] - 4), Image.LANCZOS)
    paste_x = (THUMB_SIZE[0] - page_copy.width) // 2
    paste_y = (THUMB_SIZE[1] - page_copy.height) // 2
    thumb.paste(page_copy, (paste_x, paste_y))
    return thumb


def main():
    global PPTX_FILE
    PPTX_FILE = find_pptx()
    print(f"📄 Source PPTX: {PPTX_FILE}")
    
    # 加载素材索引
    with open(INDEX_FILE, 'r') as f:
        index = json.load(f)
    
    print(f"📋 Total assets in index: {index['totalAssets']}")
    
    # 建立 slide number -> page index 映射
    unpacked_dir = PROJECT_DIR / "unpacked"
    slide_to_page = get_slide_number_mapping(unpacked_dir)
    
    if slide_to_page:
        print(f"📑 Slide to page mapping: {len(slide_to_page)} slides")
        # 检查一些映射
        sample = list(slide_to_page.items())[:5]
        print(f"   Sample: slide {sample[0][0]} -> page {sample[0][1]}, "
              f"slide {sample[-1][0]} -> page {sample[-1][1]}")
    else:
        print("⚠️  Could not build slide mapping, using slide number as page index")
    
    # 获取所有元素的坐标
    print("\n🔍 Reading element bounds from PPTX...")
    all_bounds, slide_w, slide_h, total_slides = get_element_bounds_from_pptx(PPTX_FILE)
    print(f"   {total_slides} slides, bounds for {len(all_bounds)} slides")
    
    # 按 slideNumber 分组素材
    assets_by_slide = {}
    for asset in index['assets']:
        sn = asset['sn']
        if sn not in assets_by_slide:
            assets_by_slide[sn] = []
        assets_by_slide[sn].append(asset)
    
    print(f"   Assets span {len(assets_by_slide)} unique slides")
    
    # 处理每个幻灯片
    updated = 0
    failed = 0
    total = len(index['assets'])
    
    for slide_num, slide_assets in sorted(assets_by_slide.items()):
        # 确定 page index
        if slide_to_page:
            page_idx = slide_to_page.get(slide_num)
            if page_idx is None:
                # 尝试直接用 slide_num - 1
                page_idx = slide_num - 1
        else:
            page_idx = slide_num - 1
        
        # 加载页面图片
        page_path = PAGES_DIR / f"page_{page_idx + 1}.png"
        if not page_path.exists():
            # 尝试其他命名
            page_path = PAGES_DIR / f"page_{slide_num}.png"
            if not page_path.exists():
                for a in slide_assets:
                    failed += 1
                continue
        
        page_img = Image.open(page_path)
        
        # 获取这个 slide 的元素坐标
        slide_idx = page_idx  # python-pptx 的 slide 索引 = page 索引
        bounds_list = all_bounds.get(slide_idx, [])
        
        for asset in slide_assets:
            thumb_path = PROJECT_DIR / asset['thumb']
            
            # 如果是 "slide" 类型（整页素材），直接用整页图片
            if asset.get('type') == 'slide' or len(slide_assets) == 1:
                thumb = render_full_slide_thumbnail(page_img)
                thumb.save(thumb_path, 'PNG', optimize=True)
                updated += 1
                continue
            
            # 对于 "element" 类型，需要根据元素索引找到对应的坐标
            # asset 中不存储 elementIndex 时默认为 0
            elem_idx = asset.get('ei', 0)
            
            if elem_idx > 0 and elem_idx <= len(bounds_list):
                # elementIndex 是 1-based
                bounds = bounds_list[elem_idx - 1]
                
                # 跳过非常小的元素（可能是隐藏或装饰元素）
                if bounds['width'] < 10000 or bounds['height'] < 10000:  # < ~0.01 inch
                    thumb = render_full_slide_thumbnail(page_img)
                else:
                    thumb = crop_element_thumbnail(
                        page_img, bounds, slide_w, slide_h
                    )
                
                thumb.save(thumb_path, 'PNG', optimize=True)
                updated += 1
            else:
                # 找不到对应元素，用整页图片
                thumb = render_full_slide_thumbnail(page_img)
                thumb.save(thumb_path, 'PNG', optimize=True)
                updated += 1
        
        # 进度报告
        if slide_num % 50 == 0:
            print(f"   Progress: slide {slide_num}, updated {updated}/{total}")
    
    print(f"\n{'='*50}")
    print(f"✅ Updated: {updated}")
    print(f"❌ Failed: {failed}")
    print(f"📊 Total assets: {total}")


if __name__ == "__main__":
    main()
