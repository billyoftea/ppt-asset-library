"""
extract_pages_103_191.py
从源 PPTX 中只提取第 103~191 页，生成一个新的精简 PPTX。
同时生成新的 slideIds 映射（新 PPTX 中的 slideId）。
"""

import json
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
SOURCE_PPTX = ROOT / "assets" / "source.pptx"
OUTPUT_PPTX = ROOT / "assets" / "bundle.pptx"
SLIDE_IDS_OUT = ROOT / "src" / "data" / "slideIds.json"

START_PAGE = 103  # inclusive, 1-indexed
END_PAGE = 191    # inclusive, 1-indexed

def main():
    print(f"Loading source PPTX: {SOURCE_PPTX}")
    prs = Presentation(str(SOURCE_PPTX))
    total = len(prs.slides)
    print(f"Total slides in source: {total}")
    
    if END_PAGE > total:
        print(f"Warning: END_PAGE {END_PAGE} > total {total}, clamping")
    
    # python-pptx doesn't directly support deleting slides,
    # so we'll build a new presentation by copying the XML directly.
    # Actually python-pptx supports slide deletion via low-level XML manipulation.
    
    # Get slide indices to keep (0-indexed)
    keep_indices = set(range(START_PAGE - 1, min(END_PAGE, total)))
    remove_indices = [i for i in range(total) if i not in keep_indices]
    
    print(f"Keeping slides {START_PAGE}-{min(END_PAGE, total)} ({len(keep_indices)} slides)")
    print(f"Removing {len(remove_indices)} slides")
    
    # Remove slides from end to start (to preserve indices)
    for idx in sorted(remove_indices, reverse=True):
        rId = prs.slides._sldIdLst[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[idx]
    
    remaining = len(prs.slides)
    print(f"Remaining slides: {remaining}")
    
    # Save the trimmed PPTX
    prs.save(str(OUTPUT_PPTX))
    print(f"Saved bundle PPTX: {OUTPUT_PPTX}")
    print(f"Size: {OUTPUT_PPTX.stat().st_size / 1024 / 1024:.1f} MB")
    
    # Now reload to get the new slide IDs
    prs2 = Presentation(str(OUTPUT_PPTX))
    slide_ids = {}
    for i, sld_id_entry in enumerate(prs2.slides._sldIdLst):
        original_page = START_PAGE + i  # 1-indexed original page number
        slide_id = int(sld_id_entry.get('id'))
        slide_ids[str(original_page)] = slide_id
        if i < 5 or i >= remaining - 3:
            print(f"  Page {original_page} -> slideId {slide_id}")
    
    # Save the new slideIds mapping
    with open(SLIDE_IDS_OUT, 'w') as f:
        json.dump(slide_ids, f)
    print(f"\nSaved slideIds.json with {len(slide_ids)} entries")

if __name__ == "__main__":
    main()
