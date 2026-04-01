"""
generate_hd_thumbnails.py
为每个素材生成高清插入图（从 LibreOffice 渲染的高清页面中裁剪）

- 整页素材 (slide_X)：整页缩放到 960x720
- 元素素材 (slide_X_eY)：从高清页面中按坐标裁剪，保持原始分辨率

输出到 assets/hd/ 目录，用于插入到 PPT 时使用（替代低分辨率缩略图）
"""

import json
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

# ── 配置 ──────────────────────────────────────────────
PAGES_DIR = "/tmp/pptx_render/pages"       # LibreOffice 渲染的页面 PNG
HD_OUTPUT_DIR = "assets/hd"                 # 高清图输出目录
INDEX_FILE = "src/data/assetIndex.json"
PPTX_FILE = None  # 自动检测

# 高清图最大尺寸（保持较高分辨率）
HD_MAX_SIZE = (960, 720)
# 裁剪时的边距（像素，在渲染图坐标系中）
CROP_PADDING = 8

def find_pptx():
    """找到工作目录下的 PPTX 文件"""
    import glob
    files = glob.glob("*.pptx")
    if not files:
        raise FileNotFoundError("No .pptx file found in current directory")
    return files[0]

def main():
    global PPTX_FILE
    PPTX_FILE = find_pptx()
    
    os.makedirs(HD_OUTPUT_DIR, exist_ok=True)
    
    # 加载索引
    with open(INDEX_FILE, "r") as f:
        data = json.load(f)
    
    # 加载 PPTX 获取尺寸信息
    prs = Presentation(PPTX_FILE)
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    
    # 获取渲染页面的尺寸（从第一个文件推断）
    sample_page = os.path.join(PAGES_DIR, "page_1.png")
    if not os.path.exists(sample_page):
        print(f"ERROR: Rendered pages not found in {PAGES_DIR}")
        print("Run: soffice --headless --convert-to pdf ... then render with PyMuPDF first")
        return
    
    sample_img = Image.open(sample_page)
    render_width, render_height = sample_img.size
    print(f"Slide EMU: {slide_width_emu} x {slide_height_emu}")
    print(f"Rendered page: {render_width} x {render_height}")
    
    # 计算 EMU → 像素的缩放因子
    scale_x = render_width / slide_width_emu
    scale_y = render_height / slide_height_emu
    
    # 构建每个幻灯片上元素的坐标映射
    print(f"\nBuilding element coordinate map from PPTX...")
    slide_shapes = {}  # {slide_number: [(element_index, left, top, width, height), ...]}
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        elements = []
        elem_counter = 0
        
        for shape in slide.shapes:
            # 跳过占位符（标题、页码等）
            if shape.shape_type == 14:  # PLACEHOLDER
                continue
            
            elem_counter += 1
            left = shape.left if shape.left else 0
            top = shape.top if shape.top else 0
            width = shape.width if shape.width else 0
            height = shape.height if shape.height else 0
            
            elements.append({
                "ei": elem_counter,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "name": shape.name,
            })
        
        slide_shapes[slide_num] = elements
    
    # 处理每个素材
    total = len(data["assets"])
    success = 0
    skipped = 0
    failed = 0
    
    # 缓存已加载的页面图片
    page_cache = {}
    
    for i, asset in enumerate(data["assets"]):
        asset_id = asset["id"]
        slide_num = asset["sn"]
        asset_type = asset["type"]
        elem_idx = asset.get("ei", 0)
        
        output_path = os.path.join(HD_OUTPUT_DIR, f"{asset_id}.png")
        
        # 加载对应页面
        page_path = os.path.join(PAGES_DIR, f"page_{slide_num}.png")
        if not os.path.exists(page_path):
            failed += 1
            continue
        
        if slide_num not in page_cache:
            # 清理缓存（只保留最近 5 页）
            if len(page_cache) > 5:
                oldest = list(page_cache.keys())[0]
                del page_cache[oldest]
            page_cache[slide_num] = Image.open(page_path)
        
        page_img = page_cache[slide_num]
        
        try:
            if asset_type == "slide":
                # 整页素材：直接缩放整页
                hd_img = page_img.copy()
                hd_img.thumbnail(HD_MAX_SIZE, Image.LANCZOS)
            else:
                # 元素素材：按坐标裁剪
                shapes = slide_shapes.get(slide_num, [])
                shape_info = None
                for s in shapes:
                    if s["ei"] == elem_idx:
                        shape_info = s
                        break
                
                if not shape_info:
                    # 找不到对应形状，使用整页
                    hd_img = page_img.copy()
                    hd_img.thumbnail(HD_MAX_SIZE, Image.LANCZOS)
                else:
                    # EMU → 渲染像素
                    px_left = int(shape_info["left"] * scale_x)
                    px_top = int(shape_info["top"] * scale_y)
                    px_right = px_left + int(shape_info["width"] * scale_x)
                    px_bottom = px_top + int(shape_info["height"] * scale_y)
                    
                    # 加边距
                    px_left = max(0, px_left - CROP_PADDING)
                    px_top = max(0, px_top - CROP_PADDING)
                    px_right = min(page_img.width, px_right + CROP_PADDING)
                    px_bottom = min(page_img.height, px_bottom + CROP_PADDING)
                    
                    # 确保裁剪区域有效
                    if px_right - px_left < 10 or px_bottom - px_top < 10:
                        hd_img = page_img.copy()
                        hd_img.thumbnail(HD_MAX_SIZE, Image.LANCZOS)
                    else:
                        hd_img = page_img.crop((px_left, px_top, px_right, px_bottom))
            
            # 保存为 PNG（无损）
            if hd_img.mode == "RGBA":
                bg = Image.new("RGB", hd_img.size, (255, 255, 255))
                bg.paste(hd_img, mask=hd_img.split()[3])
                hd_img = bg
            elif hd_img.mode != "RGB":
                hd_img = hd_img.convert("RGB")
            
            hd_img.save(output_path, "PNG", optimize=True)
            success += 1
            
        except Exception as e:
            failed += 1
            if failed <= 5:
                print(f"  ERROR {asset_id}: {e}")
        
        if (i + 1) % 1000 == 0:
            print(f"  Progress: {i+1}/{total} (success={success}, failed={failed})")
    
    print(f"\n✅ Done! {success} HD images generated, {failed} failed, {skipped} skipped")
    
    # 统计大小
    total_size = sum(
        os.path.getsize(os.path.join(HD_OUTPUT_DIR, f))
        for f in os.listdir(HD_OUTPUT_DIR)
        if f.endswith(".png")
    )
    print(f"Total HD images size: {total_size // (1024*1024)}MB")

if __name__ == "__main__":
    main()
