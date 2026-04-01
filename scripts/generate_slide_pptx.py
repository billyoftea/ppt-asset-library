#!/usr/bin/env python3
"""
generate_slide_pptx.py
为源 PPTX 的每张幻灯片生成一个独立的 mini PPTX 文件。
每个 mini PPTX 只包含该幻灯片（保留所有矢量形状、图片、样式），
用于通过 insertSlidesFromBase64 插入到 PowerPoint 中，实现原生可编辑插入。

输出: assets/slides/slide_{N}.pptx  (N = 幻灯片编号)
"""

import os
import sys
import json
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

PROJECT_ROOT = Path(__file__).parent.parent
SOURCE_PPTX = PROJECT_ROOT / "assets" / "source.pptx"
SLIDES_DIR = PROJECT_ROOT / "assets" / "slides"
DATA_DIR = PROJECT_ROOT / "src" / "data"


def main():
    print("=" * 60)
    print("📦 Generating per-slide mini PPTX files")
    print("=" * 60)

    if not SOURCE_PPTX.exists():
        print(f"❌ Source PPTX not found: {SOURCE_PPTX}")
        sys.exit(1)

    # 创建输出目录
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)

    # 清理旧文件
    old_count = 0
    for f in SLIDES_DIR.glob("*.pptx"):
        f.unlink()
        old_count += 1
    if old_count:
        print(f"🗑️  Cleaned {old_count} old files")

    # 加载源 PPTX
    print(f"\n📂 Loading source PPTX: {SOURCE_PPTX}")
    print(f"   File size: {os.path.getsize(SOURCE_PPTX) / 1024 / 1024:.1f} MB")
    src_prs = Presentation(str(SOURCE_PPTX))
    total_slides = len(src_prs.slides)
    print(f"   Total slides: {total_slides}")

    # 加载 assetIndex 获取哪些幻灯片有素材
    asset_index_path = DATA_DIR / "assetIndex.json"
    needed_slides = set()
    if asset_index_path.exists():
        with open(asset_index_path, 'r') as f:
            index_data = json.load(f)
        for asset in index_data.get('assets', []):
            sn = asset.get('sn', asset.get('slideNumber', 0))
            if sn > 0:
                needed_slides.add(sn)
        print(f"   Slides referenced by assets: {len(needed_slides)}")
    else:
        # 如果没有索引，为所有幻灯片生成
        needed_slides = set(range(1, total_slides + 1))
        print(f"   No asset index found, generating for all {total_slides} slides")

    generated = 0
    errors = 0
    total_size = 0

    for slide_idx, slide in enumerate(src_prs.slides):
        slide_num = slide_idx + 1

        # 只为有素材的幻灯片生成
        if slide_num not in needed_slides:
            continue

        output_path = SLIDES_DIR / f"slide_{slide_num}.pptx"

        try:
            # 创建一个新的空 PPTX，然后从源中复制这一页
            # python-pptx 不直接支持复制单页，所以我们用另一种方式：
            # 打开源 PPTX，删除所有其他幻灯片，保存为新文件
            # 但这样太慢了（每次都要加载66MB）
            #
            # 更好的方案：使用 python-pptx 的底层 lxml 操作
            generate_single_slide_pptx(src_prs, slide_idx, output_path)
            
            file_size = os.path.getsize(output_path)
            total_size += file_size
            generated += 1

            if generated % 50 == 0:
                print(f"  Progress: {generated}/{len(needed_slides)} slides generated")

        except Exception as e:
            print(f"  ⚠️  Error for slide {slide_num}: {e}")
            errors += 1

    print(f"\n{'=' * 60}")
    print(f"✅ Generation complete!")
    print(f"   Generated: {generated} mini PPTX files")
    print(f"   Errors: {errors}")
    print(f"   Total size: {total_size / 1024 / 1024:.1f} MB")
    print(f"   Average size: {total_size / max(generated, 1) / 1024:.1f} KB")
    print(f"   Output dir: {SLIDES_DIR}")


def generate_single_slide_pptx(src_prs, slide_idx, output_path):
    """
    从源 Presentation 中提取单张幻灯片，保存为独立 PPTX。
    使用 python-pptx 的底层操作来高效复制。
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from lxml import etree
    import copy as cp
    from io import BytesIO
    import zipfile

    src_slide = src_prs.slides[slide_idx]

    # 创建新的 Presentation，使用相同的幻灯片尺寸
    new_prs = Presentation()
    new_prs.slide_width = src_prs.slide_width
    new_prs.slide_height = src_prs.slide_height

    # 使用源幻灯片的 layout 对应的 slide master
    # 但由于 layout/master 不一定一样，我们用空白 layout
    blank_layout = new_prs.slide_layouts[6]  # 通常是空白 layout
    new_slide = new_prs.slides.add_slide(blank_layout)

    # 清除新幻灯片的默认内容
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'grpSp', 'pic', 'cxnSp', 'graphicFrame'):
            sp_tree.remove(child)

    # 复制源幻灯片的所有形状 XML
    src_sp_tree = src_slide.shapes._spTree
    for child in src_sp_tree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'grpSp', 'pic', 'cxnSp', 'graphicFrame'):
            new_child = cp.deepcopy(child)
            sp_tree.append(new_child)

    # 复制背景（如果有）
    src_cSld = src_slide._element
    src_bg = src_cSld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
    if src_bg is not None:
        new_cSld = new_slide._element
        new_bg = cp.deepcopy(src_bg)
        # 插入到 spTree 之前
        new_cSld_inner = new_cSld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
        if new_cSld_inner is not None:
            sp_tree_elem = new_cSld_inner.find('{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
            if sp_tree_elem is not None:
                new_cSld_inner.insert(list(new_cSld_inner).index(sp_tree_elem), new_bg)

    # 复制图片关系 - 遍历源幻灯片的所有关系
    for rel in src_slide.part.rels.values():
        if "image" in rel.reltype:
            try:
                # 获取图片的二进制数据
                image_blob = rel.target_part.blob
                content_type = rel.target_part.content_type
                partname = rel.target_part.partname

                # 在新幻灯片中创建相同的关系
                from pptx.opc.part import Part
                from pptx.opc.packuri import PackURI

                # 使用同样的 part name 以保持 rId 引用一致
                new_part = Part(
                    PackURI(str(partname)),
                    content_type,
                    image_blob,
                    new_prs.part.package
                )
                new_slide.part.rels.get_or_add(rel.rId, rel.reltype, new_part)
            except Exception:
                pass  # 某些关系可能无法复制，跳过

    new_prs.save(str(output_path))


if __name__ == '__main__':
    main()
